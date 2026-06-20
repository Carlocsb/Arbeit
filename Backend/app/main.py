from __future__ import annotations

import json
import os
import re
import sys
import textwrap
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn


def get_code_dir() -> Path:
    """
    Ordner, in dem main.py liegt.

    Entwicklung:
        Backend/app

    PyInstaller onedir:
        Backend/dist/PowerPointAutomationStudio/_internal/app
    """
    return Path(__file__).resolve().parent


def get_runtime_dir() -> Path:
    """
    Schreibbarer Runtime-Ordner.

    Entwicklung:
        Backend/runtime

    PyInstaller onedir:
        Backend/dist/PowerPointAutomationStudio/runtime
    """
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent / "runtime"

    return get_code_dir().parent / "runtime"


APP_DIR = get_code_dir()
RUNTIME_DIR = get_runtime_dir()

UPLOAD_DIR = RUNTIME_DIR / "uploads"
OUTPUT_DIR = RUNTIME_DIR / "outputs"
PREVIEW_DIR = RUNTIME_DIR / "previews"

MAPPING_DIR = APP_DIR / "mappings"

for folder in [UPLOAD_DIR, OUTPUT_DIR, PREVIEW_DIR, MAPPING_DIR]:
    folder.mkdir(parents=True, exist_ok=True)

MAPPING_PREFIX = "PPT_AUTOMATION_MAPPING::"


def safe_filename(filename: str) -> str:
    filename = Path(filename).name
    filename = re.sub(r"[^a-zA-Z0-9_.äöüÄÖÜß-]", "_", filename)
    return filename


def normalize_text(value) -> str:
    if value is None:
        return ""

    text = str(value).strip().lower()
    text = text.replace("\n", " ")
    text = re.sub(r"\s+", " ", text)
    text = text.replace("–", "-").replace("—", "-")

    return text


def normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df.columns = [str(col).strip() for col in df.columns]
    return df


def is_year_like(value) -> bool:
    try:
        year = int(float(str(value).strip()))
        return 1900 <= year <= 2100
    except Exception:
        return False


def to_number(value) -> float:
    if pd.isna(value):
        return 0.0

    if isinstance(value, (int, float)):
        return float(value)

    value = str(value).strip()
    value = value.replace(".", "")
    value = value.replace(",", ".")

    try:
        return float(value)
    except Exception:
        return 0.0


def wrap_label(label: str, width: int = 42) -> str:
    return "\n".join(textwrap.wrap(str(label), width=width))


def read_hidden_mappings_from_presentation(presentation: Presentation) -> list[dict]:
    mappings = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue

            text = shape.text_frame.text or ""

            if MAPPING_PREFIX not in text:
                continue

            raw_json = text.split(MAPPING_PREFIX, 1)[1].strip()

            try:
                mapping = json.loads(raw_json)
                mapping["_mapping_type"] = "thinkcell"
                mapping["_found_on_slide_index"] = slide_index
                mapping["_found_in_shape_index"] = shape_index
                mappings.append(mapping)
            except Exception:
                continue

    return mappings


def normalize_mapping_schema(mapping: dict) -> dict:
    normalized = dict(mapping)

    normalized.setdefault("version", "1.0")
    normalized.setdefault("mode", "thinkcell")
    normalized.setdefault("source", {})
    normalized.setdefault("layout", {})

    if not isinstance(normalized["source"], dict):
        normalized["source"] = {}

    if not isinstance(normalized["layout"], dict):
        normalized["layout"] = {}

    if "sheet" in normalized and not normalized["source"].get("sheet"):
        normalized["source"]["sheet"] = normalized.get("sheet")

    if "range" in normalized and not normalized["source"].get("range"):
        normalized["source"]["range"] = normalized.get("range")

    normalized["source"].setdefault("type", "excel_range")
    normalized["layout"].setdefault("transposed", False)
    normalized["layout"].setdefault("first_row_as_series", True)
    normalized["layout"].setdefault("first_column_as_categories", True)

    return normalized


def find_excel_range_by_columns(
    excel_path: Path,
    sheet_name: str,
    required_columns: list[str],
    header_search_rows: int = 20,
) -> str:
    if excel_path.suffix.lower() == ".xls":
        raise ValueError(
            "excel_table_columns unterstützt .xls nicht. "
            "Bitte .xlsx oder .xlsm verwenden."
        )

    workbook = load_workbook(excel_path, data_only=True)

    if sheet_name not in workbook.sheetnames:
        raise ValueError(f"Sheet '{sheet_name}' wurde nicht gefunden.")

    sheet = workbook[sheet_name]
    normalized_required = [normalize_text(col) for col in required_columns]

    for row_idx in range(1, min(sheet.max_row, header_search_rows) + 1):
        row_values = [
            sheet.cell(row=row_idx, column=col_idx).value
            for col_idx in range(1, sheet.max_column + 1)
        ]

        normalized_row = [normalize_text(value) for value in row_values]
        column_positions = []

        for required in normalized_required:
            if required not in normalized_row:
                column_positions = []
                break

            column_positions.append(normalized_row.index(required) + 1)

        if not column_positions:
            continue

        min_col = min(column_positions)
        max_col = max(column_positions)
        last_data_row = row_idx

        for data_row in range(row_idx + 1, sheet.max_row + 1):
            row_has_data = any(
                sheet.cell(row=data_row, column=col_idx).value not in [None, ""]
                for col_idx in range(min_col, max_col + 1)
            )

            if row_has_data:
                last_data_row = data_row
            else:
                break

        return (
            f"{get_column_letter(min_col)}{row_idx}:"
            f"{get_column_letter(max_col)}{last_data_row}"
        )

    raise ValueError(
        f"Keine Tabelle mit den Spalten {required_columns} "
        f"im Sheet '{sheet_name}' gefunden."
    )
def find_excel_named_table_range(
    excel_path: Path,
    table_name: str,
) -> tuple[str, str]:
    """
    Findet eine benannte Excel-Tabelle und gibt Sheet + Range zurück.
    """
    if excel_path.suffix.lower() == ".xls":
        raise ValueError(
            "excel_named_table unterstützt .xls nicht. "
            "Bitte .xlsx oder .xlsm verwenden."
        )

    workbook = load_workbook(excel_path, data_only=True)

    available_tables = []

    for worksheet in workbook.worksheets:
        for table in worksheet.tables.values():
            available_tables.append(f"{worksheet.title}.{table.name}")

            if table.name == table_name:
                return worksheet.title, table.ref

    raise ValueError(
        f"Excel-Tabelle '{table_name}' wurde nicht gefunden. "
        f"Verfügbare Tabellen: {available_tables}"
    )
def resolve_excel_range_from_mapping(excel_path: Path, mapping: dict) -> tuple[str, str]:
    """
    Unterstützt:
    - source.type = excel_range
    - source.type = excel_table_columns
    - source.type = excel_named_table
    """
    mapping = normalize_mapping_schema(mapping)
    source = mapping.get("source", {})

    source_type = source.get("type", "excel_range")

    if source_type == "excel_named_table":
        table_name = source.get("table") or mapping.get("table")

        if not table_name:
            raise ValueError("excel_named_table braucht source.table.")

        return find_excel_named_table_range(
            excel_path=excel_path,
            table_name=table_name,
        )

    sheet_name = source.get("sheet") or mapping.get("sheet")

    if not sheet_name:
        raise ValueError("Mapping braucht source.sheet.")

    if source_type == "excel_range":
        cell_range = source.get("range") or mapping.get("range")

        if not cell_range:
            raise ValueError("Mapping braucht source.range.")

        return sheet_name, cell_range

    if source_type == "excel_table_columns":
        required_columns = source.get("required_columns", [])
        header_search_rows = int(source.get("header_search_rows", 20))

        if not required_columns:
            raise ValueError("excel_table_columns braucht source.required_columns.")

        cell_range = find_excel_range_by_columns(
            excel_path=excel_path,
            sheet_name=sheet_name,
            required_columns=required_columns,
            header_search_rows=header_search_rows,
        )

        return sheet_name, cell_range

    raise ValueError(f"Unbekannter source.type: {source_type}")
def resolve_excel_range_from_mapping(excel_path: Path, mapping: dict) -> tuple[str, str]:
    mapping = normalize_mapping_schema(mapping)
    source = mapping.get("source", {})
    sheet_name = source.get("sheet") or mapping.get("sheet")

    if not sheet_name:
        raise ValueError("Mapping braucht source.sheet.")

    source_type = source.get("type", "excel_range")

    if source_type == "excel_range":
        cell_range = source.get("range") or mapping.get("range")

        if not cell_range:
            raise ValueError("Mapping braucht source.range.")

        return sheet_name, cell_range

    if source_type == "excel_table_columns":
        required_columns = source.get("required_columns", [])
        header_search_rows = int(source.get("header_search_rows", 20))

        if not required_columns:
            raise ValueError("excel_table_columns braucht source.required_columns.")

        cell_range = find_excel_range_by_columns(
            excel_path=excel_path,
            sheet_name=sheet_name,
            required_columns=required_columns,
            header_search_rows=header_search_rows,
        )

        return sheet_name, cell_range

    raise ValueError(f"Unbekannter source.type: {source_type}")


def check_thinkcell_available() -> dict:
    if os.name != "nt":
        return {
            "available": False,
            "reason": "think-cell Automation funktioniert nur unter Windows mit Excel/PowerPoint COM.",
        }

    try:
        import pythoncom
        import win32com.client
    except Exception as exc:
        return {
            "available": False,
            "reason": f"pywin32 ist nicht verfügbar: {exc}",
        }

    excel_app = None

    try:
        pythoncom.CoInitialize()

        excel_app = win32com.client.DispatchEx("Excel.Application")
        excel_app.Visible = False
        excel_app.DisplayAlerts = False

        addin = excel_app.COMAddIns("thinkcell.addin")
        _ = addin.Object

        return {
            "available": True,
            "reason": "think-cell COM API ist verfügbar.",
        }

    except Exception as exc:
        return {
            "available": False,
            "reason": f"think-cell COM API nicht verfügbar: {exc}",
        }

    finally:
        try:
            if excel_app is not None:
                excel_app.Quit()
        except Exception:
            pass

        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def update_thinkcell_charts_with_com(
    ppt_template_path: Path,
    excel_path: Path,
    output_path: Path,
    mappings: list[dict],
) -> int:
    availability = check_thinkcell_available()

    if not availability.get("available"):
        raise RuntimeError(
            availability.get("reason", "think-cell COM API ist nicht verfügbar.")
        )

    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()

    excel_app = None
    ppt_app = None
    workbook = None
    presentation = None
    updated_count = 0

    try:
        excel_app = win32com.client.DispatchEx("Excel.Application")
        ppt_app = win32com.client.DispatchEx("PowerPoint.Application")

        excel_app.Visible = False
        excel_app.DisplayAlerts = False
        ppt_app.Visible = False

        workbook = excel_app.Workbooks.Open(str(excel_path))

        presentation = ppt_app.Presentations.Open(
            str(ppt_template_path),
            WithWindow=False,
        )

        tc_addin = excel_app.COMAddIns("thinkcell.addin").Object
        tc_update = tc_addin.CreateUpdate()

        for mapping in mappings:
            mapping = normalize_mapping_schema(mapping)

            element_name = mapping.get("element_name")
            transposed = bool(mapping.get("layout", {}).get("transposed", False))

            if not element_name:
                raise ValueError("think-cell Mapping ohne element_name gefunden.")

            sheet_name, cell_range = resolve_excel_range_from_mapping(
                excel_path=excel_path,
                mapping=mapping,
            )

            try:
                worksheet = workbook.Worksheets(sheet_name)
                excel_range = worksheet.Range(cell_range)
            except Exception as error:
                raise ValueError(
                    f"Excel-Range '{sheet_name}!{cell_range}' "
                    f"für '{element_name}' ist ungültig."
                ) from error

            try:
                tc_update.AddRangeData(
                    presentation,
                    element_name,
                    excel_range,
                    transposed,
                )
            except Exception as error:
                raise RuntimeError(
                    f"think-cell konnte '{element_name}' nicht aktualisieren. "
                    "Prüfe den AddRangeData-Namen im think-cell-Diagramm."
                ) from error

            updated_count += 1

        if updated_count > 0:
            tc_update.Send()

        presentation.SaveAs(str(output_path))

        return updated_count

    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass

        try:
            if workbook is not None:
                workbook.Close(False)
        except Exception:
            pass

        try:
            if ppt_app is not None:
                ppt_app.Quit()
        except Exception:
            pass

        try:
            if excel_app is not None:
                excel_app.Quit()
        except Exception:
            pass

        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def remove_external_chart_data_link(chart):
    try:
        chart_space = chart._chartSpace
        external_data = chart_space.find(qn("c:externalData"))

        if external_data is None:
            return

        r_id = external_data.get(qn("r:id"))

        if r_id:
            try:
                chart.part.drop_rel(r_id)
            except Exception:
                pass

        chart_space.remove(external_data)

    except Exception:
        pass


def replace_or_recreate_chart(slide, shape, chart_data):
    chart = shape.chart

    try:
        remove_external_chart_data_link(chart)
        chart.replace_data(chart_data)
        return "replaced"

    except Exception as error:
        error_message = str(error)

        if (
            "target-mode is external" not in error_message
            and "_target_part" not in error_message
        ):
            raise error

        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        chart_type = chart.chart_type

        old_element = shape._element
        old_element.getparent().remove(old_element)

        slide.shapes.add_chart(
            chart_type,
            left,
            top,
            width,
            height,
            chart_data,
        )

        return "recreated"


def get_chart_categories(chart) -> list:
    try:
        return [str(category) for category in chart.plots[0].categories]
    except Exception:
        return []


def get_chart_series_names(chart) -> list:
    try:
        return [str(series.name) for series in chart.series]
    except Exception:
        return []


def find_year_column(df: pd.DataFrame):
    for col in df.columns:
        if normalize_text(col) == "jahr":
            return col

    first_col = df.columns[0]

    if df[first_col].dropna().apply(is_year_like).any():
        return first_col

    return None


def try_build_chart_data_from_wide_excel(
    df: pd.DataFrame,
    chart_categories: list,
):
    normalized_columns = {
        normalize_text(col): col
        for col in df.columns
    }

    if not any(
        normalize_text(category) in normalized_columns
        for category in chart_categories
    ):
        return None

    year_column = find_year_column(df)

    if year_column:
        first_year_value = df[year_column].dropna().iloc[0]
        series_name = (
            str(int(float(first_year_value)))
            if is_year_like(first_year_value)
            else str(first_year_value)
        )
    else:
        series_name = "Aktuelle Werte"

    categories = []
    values = []

    for category in chart_categories:
        normalized_category = normalize_text(category)

        if normalized_category in normalized_columns:
            excel_col = normalized_columns[normalized_category]
            value = df[excel_col].apply(to_number).sum()
        else:
            value = 0.0

        categories.append(category)
        values.append(float(value))

    return series_name, categories, values


def find_best_category_column(
    df: pd.DataFrame,
    chart_categories: list,
):
    normalized_chart_categories = {
        normalize_text(category)
        for category in chart_categories
        if normalize_text(category)
    }

    best_column = None
    best_score = 0

    for column in df.columns:
        values = {
            normalize_text(value)
            for value in df[column].dropna().astype(str).tolist()
        }

        score = len(normalized_chart_categories.intersection(values))

        if score > best_score:
            best_score = score
            best_column = column

    if best_score == 0:
        return None

    return best_column


def find_value_columns_for_long_excel(
    df: pd.DataFrame,
    category_column: str,
    series_names: list,
) -> list:
    excluded = {category_column}

    year_column = find_year_column(df)

    if year_column:
        excluded.add(year_column)

    for series_name in series_names:
        for col in df.columns:
            if normalize_text(col) == normalize_text(series_name):
                return [col]

    insgesamt_columns = [
        col
        for col in df.columns
        if str(col).strip().lower().startswith("insgesamt")
    ]

    if insgesamt_columns:
        return insgesamt_columns

    numeric_columns = []

    for col in df.columns:
        if col in excluded:
            continue

        converted = df[col].apply(to_number)

        if converted.sum() != 0:
            numeric_columns.append(col)

    return numeric_columns


def try_build_chart_data_from_long_excel(
    df: pd.DataFrame,
    chart_categories: list,
    series_names: list,
):
    category_column = find_best_category_column(
        df=df,
        chart_categories=chart_categories,
    )

    if not category_column:
        return None

    value_columns = find_value_columns_for_long_excel(
        df=df,
        category_column=category_column,
        series_names=series_names,
    )

    if not value_columns:
        return None

    year_column = find_year_column(df)

    if year_column:
        first_year_value = df[year_column].dropna().iloc[0]
        series_name = (
            str(int(float(first_year_value)))
            if is_year_like(first_year_value)
            else str(first_year_value)
        )
    elif series_names:
        series_name = series_names[0]
    else:
        series_name = "Aktuelle Werte"

    working_df = df.copy()
    working_df["_normalized_category"] = working_df[category_column].apply(
        normalize_text
    )

    for col in value_columns:
        working_df[col] = working_df[col].apply(to_number)

    working_df["_chart_value"] = working_df[value_columns].sum(axis=1)

    grouped = (
        working_df
        .groupby("_normalized_category")["_chart_value"]
        .sum()
        .to_dict()
    )

    categories = []
    values = []

    for category in chart_categories:
        normalized_category = normalize_text(category)
        categories.append(category)
        values.append(float(grouped.get(normalized_category, 0.0)))

    return series_name, categories, values


def replace_charts_with_matching_excel_data(
    presentation: Presentation,
    excel_path: Path,
) -> int:
    df = pd.read_excel(excel_path)

    if df.empty:
        raise ValueError("Die Excel-Datei ist leer.")

    df = normalize_columns(df)

    updated_charts = 0
    found_charts = 0
    errors = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in list(slide.shapes):
            if not getattr(shape, "has_chart", False):
                continue

            found_charts += 1

            chart = shape.chart
            chart_categories = get_chart_categories(chart)
            series_names = get_chart_series_names(chart)

            if not chart_categories:
                errors.append(
                    f"Folie {slide_index}: Diagramm hat keine auslesbaren Kategorien."
                )
                continue

            chart_data_result = try_build_chart_data_from_wide_excel(
                df=df,
                chart_categories=chart_categories,
            )

            if chart_data_result is None:
                chart_data_result = try_build_chart_data_from_long_excel(
                    df=df,
                    chart_categories=chart_categories,
                    series_names=series_names,
                )

            if chart_data_result is None:
                errors.append(
                    f"Folie {slide_index}: Keine passenden Excel-Daten "
                    "für das Diagramm gefunden."
                )
                continue

            series_name, categories, values = chart_data_result

            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(series_name, values)

            replace_or_recreate_chart(
                slide=slide,
                shape=shape,
                chart_data=chart_data,
            )

            updated_charts += 1

    if found_charts > 0 and updated_charts == 0:
        error_text = (
            "Es wurde ein Diagramm gefunden, "
            "aber es konnte nicht aktualisiert werden."
        )

        if errors:
            error_text += " Details: " + " | ".join(errors)

        raise ValueError(error_text)

    return updated_charts


def create_status_preview_images(
    presentation: Presentation,
    preview_folder: Path,
    job_id: str,
) -> list[Path]:
    preview_folder.mkdir(parents=True, exist_ok=True)

    preview_images = []
    slide_count = len(presentation.slides)

    for index in range(slide_count):
        image = Image.new(
            "RGB",
            (1920, 1080),
            color=(245, 246, 250),
        )

        draw = ImageDraw.Draw(image)

        title = f"Folie {index + 1}"
        subtitle = "PowerPoint wurde erstellt"

        try:
            font_title = ImageFont.truetype("arial.ttf", 120)
            font_subtitle = ImageFont.truetype("arial.ttf", 54)
        except Exception:
            font_title = ImageFont.load_default()
            font_subtitle = ImageFont.load_default()

        draw.text(
            (690, 430),
            title,
            fill=(44, 62, 80),
            font=font_title,
        )

        draw.text(
            (610, 570),
            subtitle,
            fill=(127, 140, 141),
            font=font_subtitle,
        )

        image_path = preview_folder / f"slide_{index + 1}.png"
        image.save(image_path)

        preview_images.append(image_path)

    return preview_images

def create_real_powerpoint_preview_images(
    pptx_path: Path,
    preview_folder: Path,
    job_id: str,
) -> list[Path]:
    """
    Exportiert echte PowerPoint-Folien als PNG.
    Funktioniert nur unter Windows mit installiertem Microsoft PowerPoint.
    """
    if os.name != "nt":
        raise RuntimeError(
            "Echte PowerPoint-Vorschau funktioniert nur unter Windows mit Microsoft PowerPoint."
        )

    try:
        import pythoncom
        import win32com.client
    except Exception as exc:
        raise RuntimeError(
            f"PowerPoint COM ist nicht verfügbar. Installiere pywin32. Details: {exc}"
        )

    pptx_path = Path(pptx_path).resolve()
    preview_folder = Path(preview_folder).resolve()
    preview_folder.mkdir(parents=True, exist_ok=True)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PowerPoint-Datei wurde nicht gefunden: {pptx_path}")

    pythoncom.CoInitialize()

    powerpoint = None
    presentation = None
    preview_images = []

    try:
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = True

        presentation = powerpoint.Presentations.Open(
            str(pptx_path),
            WithWindow=False,
            ReadOnly=True,
            Untitled=False,
        )

        slide_count = presentation.Slides.Count

        for slide_index in range(1, slide_count + 1):
            image_path = preview_folder / f"slide_{slide_index}.png"

            slide = presentation.Slides(slide_index)

            slide.Export(
                str(image_path),
                "PNG",
                1920,
                1080,
            )

            preview_images.append(image_path)

        return preview_images

    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass

        try:
            if powerpoint is not None:
                powerpoint.Quit()
        except Exception:
            pass

        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
