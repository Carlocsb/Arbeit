from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, HTMLResponse
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
from openpyxl import load_workbook
from openpyxl.utils import range_boundaries
import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import pandas as pd
import shutil
import uuid
import re
import textwrap
import sys
import os
import json
from pathlib import Path


# ------------------------------------------------------------
# EXE-sichere Pfade
# ------------------------------------------------------------

def get_app_dir() -> Path:
    """
    Gibt den Ordner zurück, in dem die App läuft.

    Normaler Python-Modus:
        Ordner der main.py

    PyInstaller-EXE-Modus:
        Ordner der .exe-Datei

    Dadurch werden uploads, outputs und previews neben der EXE erstellt.
    """
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent

    return Path(__file__).resolve().parent


def get_resource_path(relative_path: str) -> Path:
    """
    Findet gebündelte Ressourcen wie frontend.html.

    Im normalen Modus:
        neben main.py

    Im PyInstaller-Modus:
        im temporären _MEIPASS-Ordner, falls mit --add-data gebündelt
    """
    if getattr(sys, "frozen", False) and hasattr(sys, "_MEIPASS"):
        return Path(sys._MEIPASS) / relative_path

    return get_app_dir() / relative_path


APP_DIR = get_app_dir()
FRONTEND_FILE = get_resource_path("frontend.html")

UPLOAD_DIR = APP_DIR / "uploads"
OUTPUT_DIR = APP_DIR / "outputs"
PREVIEW_DIR = APP_DIR / "previews"
MAPPING_DIR = APP_DIR / "mappings"

UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)
PREVIEW_DIR.mkdir(exist_ok=True)
MAPPING_DIR.mkdir(exist_ok=True)


# ------------------------------------------------------------
# App-Konfiguration
# ------------------------------------------------------------

HOST = os.getenv("PPT_AUTOMATION_HOST", "127.0.0.1")
PORT = int(os.getenv("PPT_AUTOMATION_PORT", "8000"))
BASE_URL = os.getenv("PPT_AUTOMATION_BASE_URL", f"http://{HOST}:{PORT}")

MAPPING_PREFIX = "PPT_AUTOMATION_MAPPING::"

app = FastAPI(title="PowerPoint Automatisierung Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

app.mount("/outputs", StaticFiles(directory=OUTPUT_DIR), name="outputs")
app.mount("/previews", StaticFiles(directory=PREVIEW_DIR), name="previews")


# ------------------------------------------------------------
# Basis-Routen
# ------------------------------------------------------------

@app.get("/")
def frontend():
    if FRONTEND_FILE.exists():
        return FileResponse(FRONTEND_FILE)

    return HTMLResponse(
        """
        <html>
            <head>
                <title>PowerPoint Automation Backend</title>
            </head>
            <body style="font-family:Arial;padding:40px;">
                <h1>PowerPoint Automation Backend läuft</h1>
                <p>Die Datei <strong>frontend.html</strong> wurde nicht gefunden.</p>
                <p>Lege frontend.html neben die EXE oder bündele sie mit PyInstaller per --add-data.</p>
            </body>
        </html>
        """
    )


@app.get("/health")
def health_check():
    return {
        "status": "online",
        "app_dir": str(APP_DIR),
        "upload_dir": str(UPLOAD_DIR),
        "output_dir": str(OUTPUT_DIR),
        "preview_dir": str(PREVIEW_DIR),
        "mapping_dir": str(MAPPING_DIR),
        "frontend_file": str(FRONTEND_FILE),
        "frontend_exists": FRONTEND_FILE.exists(),
        "base_url": BASE_URL,
    }



@app.get("/app-info")
def app_info():
    """
    Kleine Info-Route für Desktop/WebView-Frontend.
    """
    return {
        "app_name": "PowerPoint Automation Studio",
        "mode": "desktop_webview_ready",
        "base_url": BASE_URL,
        "frontend_exists": FRONTEND_FILE.exists(),
    }


@app.get("/thinkcell-health")
def thinkcell_health():
    """
    Prüft lokal, ob die think-cell COM API verfügbar ist.

    Diese Route funktioniert nur auf Windows-Rechnern mit installierter
    Desktop-Version von Excel, PowerPoint und think-cell.
    """
    return check_thinkcell_available()


# ------------------------------------------------------------
# Standard-Generierung: Excel + PPTX
# ------------------------------------------------------------

@app.post("/generate-ppt")
async def generate_ppt(
    excel: UploadFile = File(...),
    template: UploadFile = File(...)
):
    if not excel.filename.endswith((".xlsx", ".xls")):
        raise HTTPException(status_code=400, detail="Bitte eine Excel-Datei hochladen.")

    if not template.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Bitte eine PowerPoint-Datei hochladen.")

    job_id = str(uuid.uuid4())

    excel_path = UPLOAD_DIR / f"{job_id}_{safe_filename(excel.filename)}"
    template_path = UPLOAD_DIR / f"{job_id}_{safe_filename(template.filename)}"
    output_path = OUTPUT_DIR / f"{job_id}_Aktualisierte_Praesentation.pptx"

    save_upload_file(excel, excel_path)
    save_upload_file(template, template_path)

    try:
        presentation = Presentation(template_path)

        # Optional: Text-Platzhalter ersetzen, falls die Excel Placeholder/Value enthält.
        data = read_excel_mapping(excel_path)
        replace_placeholders_in_presentation(
            presentation=presentation,
            data=data
        )

        # PowerPoint-Diagramme anhand der neuen Excel-Daten aktualisieren.
        updated_charts = replace_charts_with_matching_excel_data(
            presentation=presentation,
            excel_path=excel_path
        )

        presentation.save(output_path)

        preview_folder = PREVIEW_DIR / job_id
        preview_folder.mkdir(exist_ok=True)

        preview_images = create_excel_chart_preview_images(
            excel_path=excel_path,
            preview_folder=preview_folder,
            job_id=job_id
        )

        if not preview_images:
            preview_images = create_status_preview_images(
                presentation=presentation,
                preview_folder=preview_folder,
                job_id=job_id
            )

        return {
            "message": "PowerPoint erfolgreich erstellt.",
            "updated_charts": updated_charts,
            "ppt_file": f"{BASE_URL}/outputs/{output_path.name}",
            "preview_images": preview_images
        }

    except Exception as error:
        raise HTTPException(status_code=500, detail=str(error))


# ------------------------------------------------------------
# Neue Routen: PPTX / Excel analysieren
# ------------------------------------------------------------

@app.post("/inspect-ppt")
async def inspect_ppt(
    template: UploadFile = File(...)
):
    """
    Liest eine PPTX aus und gibt Folien, native Charts und versteckte Mappings zurück.
    """
    if not template.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Bitte eine PowerPoint-Datei hochladen.")

    job_id = str(uuid.uuid4())
    template_path = UPLOAD_DIR / f"{job_id}_{safe_filename(template.filename)}"

    save_upload_file(template, template_path)

    try:
        presentation = Presentation(template_path)

        return {
            "message": "PowerPoint erfolgreich analysiert.",
            "slides": inspect_presentation(presentation),
            "hidden_mappings": read_hidden_mappings_from_presentation(presentation)
        }

    except Exception as error:
        raise HTTPException(status_code=500, detail=str(error))


@app.post("/inspect-excel")
async def inspect_excel(
    excel: UploadFile = File(...),
    preview_rows: int = Form(12),
    preview_cols: int = Form(10)
):
    """
    Liest Excel-Sheets aus und zeigt Used-Range plus kleine Vorschau.
    Damit kann das Frontend Sheet und Range auswählen.
    """
    if not excel.filename.endswith((".xlsx", ".xlsm", ".xls")):
        raise HTTPException(status_code=400, detail="Bitte eine Excel-Datei hochladen.")

    job_id = str(uuid.uuid4())
    excel_path = UPLOAD_DIR / f"{job_id}_{safe_filename(excel.filename)}"

    save_upload_file(excel, excel_path)

    try:
        workbook_info = inspect_excel_workbook(
            excel_path=excel_path,
            preview_rows=preview_rows,
            preview_cols=preview_cols
        )

        return {
            "message": "Excel erfolgreich analysiert.",
            "excel_file": str(excel_path.name),
            "sheets": workbook_info
        }

    except Exception as error:
        raise HTTPException(status_code=500, detail=str(error))


# ------------------------------------------------------------
# Neue Routen: Mapping versteckt in PPTX speichern / lesen
# ------------------------------------------------------------

@app.post("/write-hidden-mappings")
async def write_hidden_mappings(
    template: UploadFile = File(...),
    mappings_json: str = Form(...)
):
    """
    Schreibt Mapping-Infos versteckt in die jeweiligen Folien.

    mappings_json Beispiel:
    [
        {
            "element_name": "umsatz_chart",
            "slide_index": 3,
            "shape_index": 17,
            "mode": "thinkcell",
            "chart_type": "column",
            "source": {
                "type": "excel_range",
                "sheet": "Umsatz",
                "range": "A1:D4"
            },
            "layout": {
                "first_row_as_series": true,
                "first_column_as_categories": true
            }
        }
    ]
    """
    if not template.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Bitte eine PowerPoint-Datei hochladen.")

    try:
        mappings = json.loads(mappings_json)
        if not isinstance(mappings, list):
            raise ValueError("mappings_json muss eine Liste sein.")
    except Exception as error:
        raise HTTPException(status_code=400, detail=f"Ungültiges mappings_json: {error}")

    job_id = str(uuid.uuid4())
    template_path = UPLOAD_DIR / f"{job_id}_{safe_filename(template.filename)}"
    output_path = OUTPUT_DIR / f"{job_id}_Template_mit_Mapping.pptx"
    mapping_file_path = MAPPING_DIR / f"{job_id}_mapping.json"

    save_upload_file(template, template_path)

    try:
        presentation = Presentation(template_path)

        written_count = write_hidden_mappings_to_presentation(
            presentation=presentation,
            mappings=mappings
        )

        presentation.save(output_path)

        with mapping_file_path.open("w", encoding="utf-8") as file:
            json.dump(mappings, file, ensure_ascii=False, indent=2)

        return {
            "message": "Mappings wurden versteckt in die PowerPoint geschrieben.",
            "written_mappings": written_count,
            "ppt_file": f"{BASE_URL}/outputs/{output_path.name}",
            "mapping_file": str(mapping_file_path.name),
            "mappings": mappings
        }

    except Exception as error:
        raise HTTPException(status_code=500, detail=str(error))


@app.post("/read-hidden-mappings")
async def read_hidden_mappings(
    template: UploadFile = File(...)
):
    """
    Liest versteckte Mappings aus einer PPTX.
    """
    if not template.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Bitte eine PowerPoint-Datei hochladen.")

    job_id = str(uuid.uuid4())
    template_path = UPLOAD_DIR / f"{job_id}_{safe_filename(template.filename)}"

    save_upload_file(template, template_path)

    try:
        presentation = Presentation(template_path)
        mappings = read_hidden_mappings_from_presentation(presentation)

        return {
            "message": "Versteckte Mappings wurden gelesen.",
            "mapping_count": len(mappings),
            "mappings": mappings
        }

    except Exception as error:
        raise HTTPException(status_code=500, detail=str(error))


@app.post("/generate-ppt-from-hidden-mappings")
async def generate_ppt_from_hidden_mappings(
    excel: UploadFile = File(...),
    template: UploadFile = File(...)
):
    """
    Liest versteckte Mappings aus der PowerPoint und aktualisiert damit:

    - native PowerPoint-Charts per python-pptx
    - think-cell-Charts per Windows COM API + think-cell AddRangeData

    Unterstützte Mapping-Formate:

    Kompakt/flach:
        PPT_AUTOMATION_MAPPING::{
          "mode":"thinkcell",
          "element_name":"umsatz_chart",
          "slide_index":3,
          "chart_type":"column",
          "sheet":"Umsatz",
          "range":"A1:D4"
        }

    Empfohlen/strukturiert:
        PPT_AUTOMATION_MAPPING::{
          "version":"1.0",
          "mode":"thinkcell",
          "element_name":"umsatz_chart",
          "slide_index":3,
          "chart_type":"column",
          "source":{"type":"excel_range","sheet":"Umsatz","range":"A1:D4"},
          "layout":{"transposed":false}
        }
    """
    if not excel.filename.endswith((".xlsx", ".xlsm", ".xls")):
        raise HTTPException(status_code=400, detail="Bitte eine Excel-Datei hochladen.")

    if not template.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Bitte eine PowerPoint-Datei hochladen.")

    job_id = str(uuid.uuid4())

    excel_path = UPLOAD_DIR / f"{job_id}_{safe_filename(excel.filename)}"
    template_path = UPLOAD_DIR / f"{job_id}_{safe_filename(template.filename)}"
    temp_native_path = OUTPUT_DIR / f"{job_id}_temp_native.pptx"
    output_path = OUTPUT_DIR / f"{job_id}_ThinkCell_Aktualisiert.pptx"

    save_upload_file(excel, excel_path)
    save_upload_file(template, template_path)

    try:
        presentation = Presentation(template_path)
        raw_mappings = read_hidden_mappings_from_presentation(presentation)
        mappings = [normalize_mapping_schema(mapping) for mapping in raw_mappings]

        if not mappings:
            raise ValueError("Keine versteckten Mappings in der PowerPoint gefunden.")

        native_mappings = []
        thinkcell_mappings = []
        unsupported_mappings = []

        for mapping in mappings:
            mode = normalize_text(mapping.get("mode", ""))

            if mode in ["native", "native_powerpoint", "powerpoint", "ppt_chart"]:
                native_mappings.append(mapping)

            elif mode in ["thinkcell", "think-cell", "think_cell"]:
                thinkcell_mappings.append(mapping)

            else:
                unsupported_mappings.append({
                    "element_name": mapping.get("element_name"),
                    "slide_index": mapping.get("slide_index"),
                    "reason": f"Unbekannter Mapping-Modus: {mapping.get('mode')}"
                })

        updated_native_charts = 0
        updated_thinkcell_charts = 0

        # 1) Native PowerPoint-Charts zuerst aktualisieren.
        for mapping in native_mappings:
            update_native_chart_from_mapping(
                presentation=presentation,
                excel_path=excel_path,
                mapping=mapping
            )
            updated_native_charts += 1

        # Nur speichern, wenn python-pptx wirklich Änderungen vorgenommen hat.
        # So bleibt die Original-PPTX für reine think-cell-Templates möglichst unverändert.
        if native_mappings:
            presentation.save(temp_native_path)
            com_input_path = temp_native_path
        else:
            com_input_path = template_path

        # 2) think-cell-Charts per COM API aktualisieren.
        if thinkcell_mappings:
            updated_thinkcell_charts = update_thinkcell_charts_with_com(
                ppt_template_path=com_input_path,
                excel_path=excel_path,
                output_path=output_path,
                mappings=thinkcell_mappings
            )
        else:
            # Falls nur native Charts vorhanden waren, die temporäre Datei als Output nutzen.
            if native_mappings:
                shutil.copyfile(temp_native_path, output_path)
            else:
                shutil.copyfile(template_path, output_path)

        preview_folder = PREVIEW_DIR / job_id
        preview_folder.mkdir(exist_ok=True)

        preview_images = create_status_preview_images(
            presentation=Presentation(output_path),
            preview_folder=preview_folder,
            job_id=job_id
        )

        return {
            "message": "PowerPoint wurde anhand der versteckten Mappings aktualisiert.",
            "updated_native_charts": updated_native_charts,
            "updated_thinkcell_charts": updated_thinkcell_charts,
            "unsupported_mappings": unsupported_mappings,
            "ppt_file": f"{BASE_URL}/outputs/{output_path.name}",
            "preview_images": preview_images,
            "mappings": mappings
        }

    except Exception as error:
        raise HTTPException(status_code=500, detail=str(error))



# ------------------------------------------------------------
# think-cell Automation über Windows COM
# ------------------------------------------------------------

def check_thinkcell_available() -> dict:
    """
    Prüft, ob Excel, PowerPoint und das think-cell COM Add-in lokal verfügbar sind.
    Diese Prüfung startet Excel kurz im Hintergrund und beendet es danach wieder.
    """
    if os.name != "nt":
        return {
            "available": False,
            "reason": "think-cell Automation funktioniert nur unter Windows."
        }

    try:
        import pythoncom
        import win32com.client
    except Exception:
        return {
            "available": False,
            "reason": "pywin32 ist nicht installiert. Ergänze pywin32 in requirements.txt."
        }

    pythoncom.CoInitialize()
    excel_app = None

    try:
        excel_app = win32com.client.DispatchEx("Excel.Application")
        excel_app.Visible = False
        excel_app.DisplayAlerts = False

        try:
            addin = excel_app.COMAddIns("thinkcell.addin")
        except Exception:
            return {
                "available": False,
                "reason": "think-cell Add-in wurde in Excel nicht gefunden."
            }

        try:
            _ = addin.Object
        except Exception:
            return {
                "available": False,
                "reason": "think-cell Add-in ist vorhanden, aber nicht aktiv oder nicht lizenziert."
            }

        return {
            "available": True,
            "reason": "think-cell COM API ist verfügbar."
        }

    except Exception as error:
        return {
            "available": False,
            "reason": str(error)
        }

    finally:
        try:
            if excel_app is not None:
                excel_app.Quit()
        except Exception:
            pass

        pythoncom.CoUninitialize()


def normalize_mapping_schema(mapping: dict) -> dict:
    """
    Vereinheitlicht flache und strukturierte Mapping-Formate.

    Unterstützt z. B.:
        {"element_name":"umsatz_chart", "sheet":"Umsatz", "range":"A1:D4"}

    und wandelt es intern in:
        {"source":{"sheet":"Umsatz", "range":"A1:D4"}, "layout":{...}}
    """
    normalized = dict(mapping)

    normalized.setdefault("version", "1.0")
    normalized.setdefault("mode", "thinkcell")
    normalized.setdefault("source", {})
    normalized.setdefault("layout", {})

    if not isinstance(normalized["source"], dict):
        normalized["source"] = {}

    if not isinstance(normalized["layout"], dict):
        normalized["layout"] = {}

    # Flache Felder aus älteren/einfacheren Mappings übernehmen.
    if "sheet" in normalized and not normalized["source"].get("sheet"):
        normalized["source"]["sheet"] = normalized.get("sheet")

    if "range" in normalized and not normalized["source"].get("range"):
        normalized["source"]["range"] = normalized.get("range")

    normalized["source"].setdefault("type", "excel_range")
    normalized["layout"].setdefault("transposed", False)
    normalized["layout"].setdefault("first_row_as_series", True)
    normalized["layout"].setdefault("first_column_as_categories", True)

    return normalized


def update_thinkcell_charts_with_com(
    ppt_template_path: Path,
    excel_path: Path,
    output_path: Path,
    mappings: list[dict]
) -> int:
    """
    Aktualisiert think-cell-Diagramme über PowerPoint + Excel + think-cell COM API.

    Voraussetzungen beim Nutzer:
    - Windows
    - Microsoft Excel Desktop
    - Microsoft PowerPoint Desktop
    - think-cell installiert, aktiviert und lizenziert
    - think-cell Elemente haben AddRangeData Names, die zu element_name passen
    """
    availability = check_thinkcell_available()

    if not availability.get("available"):
        raise RuntimeError(availability.get("reason", "think-cell COM API ist nicht verfügbar."))

    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()

    excel_app = None
    ppt_app = None
    workbook = None
    presentation = None
    updated_count = 0

    try:
        excel_app = win32com.client.DispatchEx("Excel.Application")
        ppt_app = win32com.client.DispatchEx("PowerPoint.Application")

        excel_app.Visible = False
        excel_app.DisplayAlerts = False
        ppt_app.Visible = False

        workbook = excel_app.Workbooks.Open(str(excel_path))
        presentation = ppt_app.Presentations.Open(
            str(ppt_template_path),
            WithWindow=False
        )

        try:
            tc_addin = excel_app.COMAddIns("thinkcell.addin").Object
        except Exception as error:
            raise RuntimeError(
                "think-cell COM Add-in wurde in Excel nicht gefunden oder ist nicht aktiv."
            ) from error

        tc_update = tc_addin.CreateUpdate()

        for mapping in mappings:
            mapping = normalize_mapping_schema(mapping)
            mode = normalize_text(mapping.get("mode", ""))

            if mode not in ["thinkcell", "think-cell", "think_cell"]:
                continue

            element_name = mapping.get("element_name")
            source = mapping.get("source", {})
            sheet_name = source.get("sheet")
            cell_range = source.get("range")
            transposed = bool(mapping.get("layout", {}).get("transposed", False))

            if not element_name:
                raise ValueError("think-cell Mapping ohne element_name gefunden.")

            if not sheet_name or not cell_range:
                raise ValueError(
                    f"Mapping für '{element_name}' benötigt source.sheet und source.range."
                )

            try:
                worksheet = workbook.Worksheets(sheet_name)
            except Exception as error:
                raise ValueError(
                    f"Excel-Sheet '{sheet_name}' für Mapping '{element_name}' wurde nicht gefunden."
                ) from error

            try:
                excel_range = worksheet.Range(cell_range)
            except Exception as error:
                raise ValueError(
                    f"Excel-Range '{sheet_name}!{cell_range}' für Mapping '{element_name}' ist ungültig."
                ) from error

            try:
                tc_update.AddRangeData(
                    presentation,
                    element_name,
                    excel_range,
                    transposed
                )
            except Exception as error:
                raise RuntimeError(
                    f"think-cell konnte das Element '{element_name}' nicht aktualisieren. "
                    "Prüfe, ob der AddRangeData Name im think-cell-Diagramm exakt gleich gesetzt ist."
                ) from error

            updated_count += 1

        if updated_count > 0:
            tc_update.Send()

        presentation.SaveAs(str(output_path))
        return updated_count

    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass

        try:
            if workbook is not None:
                workbook.Close(False)
        except Exception:
            pass

        try:
            if ppt_app is not None:
                ppt_app.Quit()
        except Exception:
            pass

        try:
            if excel_app is not None:
                excel_app.Quit()
        except Exception:
            pass

        pythoncom.CoUninitialize()


# ------------------------------------------------------------
# Upload- und Dateihilfen
# ------------------------------------------------------------

def safe_filename(filename: str) -> str:
    """
    Entfernt problematische Zeichen aus Dateinamen.
    """
    filename = Path(filename).name
    filename = re.sub(r"[^a-zA-Z0-9_.äöüÄÖÜß-]", "_", filename)
    return filename


def save_upload_file(upload_file: UploadFile, destination: Path):
    with destination.open("wb") as buffer:
        shutil.copyfileobj(upload_file.file, buffer)


# ------------------------------------------------------------
# Excel-Analyse für Mapping-Assistent
# ------------------------------------------------------------

def inspect_excel_workbook(
    excel_path: Path,
    preview_rows: int = 12,
    preview_cols: int = 10
) -> list:
    workbook = load_workbook(excel_path, data_only=True)
    result = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]

        max_row = sheet.max_row or 1
        max_col = sheet.max_column or 1

        used_range = f"A1:{column_number_to_letter(max_col)}{max_row}"

        preview = []

        for row in sheet.iter_rows(
            min_row=1,
            max_row=min(max_row, preview_rows),
            min_col=1,
            max_col=min(max_col, preview_cols),
            values_only=True
        ):
            preview.append([
                "" if value is None else value
                for value in row
            ])

        result.append({
            "name": sheet_name,
            "used_range": used_range,
            "max_row": max_row,
            "max_column": max_col,
            "preview": preview
        })

    return result


def column_number_to_letter(column_number: int) -> str:
    result = ""

    while column_number:
        column_number, remainder = divmod(column_number - 1, 26)
        result = chr(65 + remainder) + result

    return result


def read_excel_range_as_table(
    excel_path: Path,
    sheet_name: str,
    cell_range: str
) -> list:
    workbook = load_workbook(excel_path, data_only=True)

    if sheet_name not in workbook.sheetnames:
        raise ValueError(f"Sheet '{sheet_name}' wurde in der Excel-Datei nicht gefunden.")

    sheet = workbook[sheet_name]

    min_col, min_row, max_col, max_row = range_boundaries(cell_range)

    table = []

    for row in sheet.iter_rows(
        min_row=min_row,
        max_row=max_row,
        min_col=min_col,
        max_col=max_col,
        values_only=True
    ):
        table.append([
            "" if value is None else value
            for value in row
        ])

    return table


# ------------------------------------------------------------
# Versteckte Mapping-Felder in PowerPoint
# ------------------------------------------------------------

def write_hidden_mappings_to_presentation(
    presentation: Presentation,
    mappings: list[dict]
) -> int:
    written_count = 0

    for mapping in mappings:
        slide_index = int(mapping.get("slide_index", 0))

        if slide_index < 1 or slide_index > len(presentation.slides):
            raise ValueError(f"Ungültige slide_index im Mapping: {slide_index}")

        slide = presentation.slides[slide_index - 1]

        element_name = mapping.get("element_name", f"slide_{slide_index}_mapping")
        mapping = enrich_mapping_defaults(mapping)

        remove_existing_hidden_mapping_from_slide(
            slide=slide,
            element_name=element_name
        )

        add_hidden_mapping_to_slide(
            slide=slide,
            mapping=mapping
        )

        written_count += 1

    return written_count


def enrich_mapping_defaults(mapping: dict) -> dict:
    """
    Ergänzt Mapping-Defaults und unterstützt auch flache Felder wie sheet/range.
    """
    return normalize_mapping_schema(mapping)


def add_hidden_mapping_to_slide(slide, mapping: dict):
    """
    Schreibt ein technisches Mapping unsichtbar in eine Folie.
    Das Mapping kann später mit python-pptx wieder ausgelesen werden.

    Strategie:
    - Textfeld außerhalb der sichtbaren Folie
    - Schriftgröße 1
    - weiße Schrift
    """
    mapping_text = MAPPING_PREFIX + json.dumps(mapping, ensure_ascii=False)

    textbox = slide.shapes.add_textbox(
        Inches(-1.5),
        Inches(-1.5),
        Inches(0.5),
        Inches(0.2)
    )

    textbox.name = f"ppt_automation_mapping_{mapping.get('element_name', 'unknown')}"

    text_frame = textbox.text_frame
    text_frame.clear()

    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = mapping_text

    font = run.font
    font.size = Pt(1)
    font.color.rgb = RGBColor(255, 255, 255)

    return textbox


def remove_existing_hidden_mapping_from_slide(slide, element_name: str):
    """
    Entfernt vorhandene Mapping-Felder für denselben element_name,
    damit beim erneuten Speichern keine Duplikate entstehen.
    """
    shapes = list(slide.shapes)

    for shape in shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text_frame.text or ""

        if MAPPING_PREFIX not in text:
            continue

        try:
            raw_json = text.split(MAPPING_PREFIX, 1)[1].strip()
            existing_mapping = json.loads(raw_json)

            if existing_mapping.get("element_name") == element_name:
                old_element = shape._element
                old_element.getparent().remove(old_element)

        except Exception:
            continue


def read_hidden_mappings_from_presentation(presentation: Presentation) -> list:
    """
    Liest alle versteckten PPT_AUTOMATION_MAPPING-Felder aus der Präsentation.
    """
    mappings = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):

            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text or ""

            if MAPPING_PREFIX not in text:
                continue

            raw_json = text.split(MAPPING_PREFIX, 1)[1].strip()

            try:
                mapping = json.loads(raw_json)
                mapping["_found_on_slide_index"] = slide_index
                mapping["_found_in_shape_index"] = shape_index
                mappings.append(mapping)
            except Exception:
                continue

    return mappings


# ------------------------------------------------------------
# PowerPoint-Inspektion
# ------------------------------------------------------------

def inspect_presentation(presentation: Presentation) -> list:
    slides_result = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_data = {
            "slide_index": slide_index,
            "texts": [],
            "tables": [],
            "native_charts": [],
            "possible_thinkcell_shapes": [],
            "hidden_mappings": []
        }

        native_chart_counter = 0

        for shape_index, shape in enumerate(slide.shapes, start=1):

            if shape.has_text_frame:
                text = shape.text_frame.text or ""

                if MAPPING_PREFIX in text:
                    try:
                        raw_json = text.split(MAPPING_PREFIX, 1)[1].strip()
                        mapping = json.loads(raw_json)
                        mapping["_shape_index"] = shape_index
                        slide_data["hidden_mappings"].append(mapping)
                    except Exception:
                        pass
                elif text.strip():
                    slide_data["texts"].append({
                        "shape_index": shape_index,
                        "text": text
                    })

            if shape.has_table:
                rows = []

                for row in shape.table.rows:
                    rows.append([
                        cell.text for cell in row.cells
                    ])

                slide_data["tables"].append({
                    "shape_index": shape_index,
                    "rows": rows
                })

            if shape.has_chart:
                native_chart_counter += 1
                chart = shape.chart

                slide_data["native_charts"].append({
                    "chart_index": native_chart_counter,
                    "shape_index": shape_index,
                    "shape_name": getattr(shape, "name", ""),
                    "chart_type": str(chart.chart_type),
                    "categories": get_chart_categories(chart),
                    "series_names": get_chart_series_names(chart)
                })

            if is_possible_thinkcell_shape(shape):
                slide_data["possible_thinkcell_shapes"].append({
                    "shape_index": shape_index,
                    "shape_name": getattr(shape, "name", ""),
                    "shape_type": str(getattr(shape, "shape_type", "")),
                    "left": int(getattr(shape, "left", 0)),
                    "top": int(getattr(shape, "top", 0)),
                    "width": int(getattr(shape, "width", 0)),
                    "height": int(getattr(shape, "height", 0)),
                    "suggested_element_name": f"slide_{slide_index}_shape_{shape_index}"
                })

        slides_result.append(slide_data)

    return slides_result


def is_possible_thinkcell_shape(shape) -> bool:
    """
    Einfache Heuristik für mögliche think-cell-Elemente.

    Hinweis:
    Das ist keine perfekte think-cell-Erkennung.
    Die sichere Aktualisierung kommt später über PowerPoint + think-cell Automation.
    """
    if shape.has_chart:
        return False

    shape_name = normalize_text(getattr(shape, "name", ""))
    xml_text = ""

    try:
        xml_text = normalize_text(shape._element.xml)
    except Exception:
        xml_text = ""

    keywords = [
        "think-cell",
        "thinkcell",
        "think cell",
        "tc_",
        "ppttc"
    ]

    if any(keyword in shape_name for keyword in keywords):
        return True

    if any(keyword in xml_text for keyword in keywords):
        return True

    return False


# ------------------------------------------------------------
# Native Charts aus verstecktem Mapping aktualisieren
# ------------------------------------------------------------

def update_native_chart_from_mapping(
    presentation: Presentation,
    excel_path: Path,
    mapping: dict
):
    slide_index = int(mapping.get("slide_index", 0))
    chart_index = int(mapping.get("chart_index", 0))

    if slide_index < 1 or slide_index > len(presentation.slides):
        raise ValueError(f"Ungültige slide_index im Mapping: {slide_index}")

    if chart_index < 1:
        raise ValueError("Für native Charts muss chart_index im Mapping gesetzt sein.")

    source = mapping.get("source", {})
    sheet_name = source.get("sheet")
    cell_range = source.get("range")

    if not sheet_name or not cell_range:
        raise ValueError("Mapping benötigt source.sheet und source.range.")

    table = read_excel_range_as_table(
        excel_path=excel_path,
        sheet_name=sheet_name,
        cell_range=cell_range
    )

    layout = mapping.get("layout", {})
    chart_data = build_category_chart_data_from_table(
        table=table,
        first_row_as_series=bool(layout.get("first_row_as_series", True)),
        first_column_as_categories=bool(layout.get("first_column_as_categories", True))
    )

    slide = presentation.slides[slide_index - 1]

    current_chart_index = 0

    for shape in slide.shapes:
        if not shape.has_chart:
            continue

        current_chart_index += 1

        if current_chart_index == chart_index:
            replace_or_recreate_chart(
                slide=slide,
                shape=shape,
                chart_data=chart_data
            )
            return

    raise ValueError(f"Natives Diagramm {chart_index} auf Folie {slide_index} wurde nicht gefunden.")


def build_category_chart_data_from_table(
    table: list[list],
    first_row_as_series: bool = True,
    first_column_as_categories: bool = True
) -> CategoryChartData:
    """
    Baut CategoryChartData aus einem Excel-Zellbereich.

    Standardlayout:
        A1 leer / Header
        B1:D1 = Reihen
        A2:A4 = Kategorien
        B2:D4 = Werte
    """
    if not table or not table[0]:
        raise ValueError("Der Excel-Bereich ist leer.")

    chart_data = CategoryChartData()

    if first_row_as_series and first_column_as_categories:
        categories = [str(row[0]) for row in table[1:]]

        chart_data.categories = categories

        series_names = [str(value) for value in table[0][1:]]

        for col_offset, series_name in enumerate(series_names, start=1):
            values = []

            for row in table[1:]:
                value = row[col_offset] if col_offset < len(row) else 0
                values.append(to_number(value))

            chart_data.add_series(series_name, values)

        return chart_data

    if not first_row_as_series and first_column_as_categories:
        categories = [str(row[0]) for row in table]
        values = [to_number(row[1]) if len(row) > 1 else 0 for row in table]

        chart_data.categories = categories
        chart_data.add_series("Werte", values)

        return chart_data

    raise ValueError("Dieses Tabellenlayout wird aktuell noch nicht unterstützt.")


# ------------------------------------------------------------
# PowerPoint-Diagramm-Hilfsfunktionen
# ------------------------------------------------------------

def remove_external_chart_data_link(chart):
    """
    Entfernt externe Excel-Verknüpfungen aus einem PowerPoint-Diagramm.
    """
    try:
        chart_space = chart._chartSpace
        external_data = chart_space.find(qn("c:externalData"))

        if external_data is None:
            return

        r_id = external_data.get(qn("r:id"))

        if r_id:
            try:
                chart.part.drop_rel(r_id)
            except Exception:
                pass

        chart_space.remove(external_data)

    except Exception:
        pass


def replace_or_recreate_chart(slide, shape, chart_data):
    """
    Versucht zuerst, die Daten des bestehenden Diagramms zu ersetzen.
    Wenn das wegen einer externen Excel-Verknüpfung fehlschlägt,
    wird das Diagramm gelöscht und an derselben Position neu erstellt.
    """
    chart = shape.chart

    try:
        remove_external_chart_data_link(chart)
        chart.replace_data(chart_data)
        return "replaced"

    except Exception as error:
        error_message = str(error)

        if "target-mode is external" not in error_message and "_target_part" not in error_message:
            raise error

        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        chart_type = chart.chart_type

        old_element = shape._element
        old_element.getparent().remove(old_element)

        slide.shapes.add_chart(
            chart_type,
            left,
            top,
            width,
            height,
            chart_data
        )

        return "recreated"


def normalize_text(value) -> str:
    """
    Macht Texte vergleichbar.
    """
    if value is None:
        return ""

    text = str(value).strip().lower()
    text = text.replace("\n", " ")
    text = re.sub(r"\s+", " ", text)
    text = text.replace("–", "-").replace("—", "-")
    return text


def normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    df.columns = [str(col).strip() for col in df.columns]
    return df


def is_year_like(value) -> bool:
    try:
        year = int(float(str(value).strip()))
        return 1900 <= year <= 2100
    except Exception:
        return False


def to_number(value) -> float:
    """
    Wandelt Excel-Werte robust in Zahlen um.
    """
    if pd.isna(value):
        return 0.0

    if isinstance(value, (int, float)):
        return float(value)

    value = str(value).strip()
    value = value.replace(".", "")
    value = value.replace(",", ".")

    try:
        return float(value)
    except Exception:
        return 0.0


def wrap_label(label: str, width: int = 42) -> str:
    """
    Kürzt lange Achsenbeschriftungen optisch durch Zeilenumbrüche.
    """
    label = str(label)
    return "\n".join(textwrap.wrap(label, width=width))


# ------------------------------------------------------------
# Text-Platzhalter
# ------------------------------------------------------------

def read_excel_mapping(excel_path: Path) -> dict:
    """
    Liest nur dann Text-Platzhalter ein, wenn die Excel-Datei wirklich
    eine Placeholder/Value-Struktur hat.

    Beispiel:
    Placeholder | Value
    TITEL       | Analyse 2017
    """
    df = pd.read_excel(excel_path)

    if df.empty:
        return {}

    df = normalize_columns(df)

    columns_lower = [str(col).lower().strip() for col in df.columns]

    if "placeholder" not in columns_lower or "value" not in columns_lower:
        return {}

    placeholder_col = df.columns[columns_lower.index("placeholder")]
    value_col = df.columns[columns_lower.index("value")]

    data = {}

    for _, row in df.iterrows():
        key = str(row[placeholder_col]).strip()
        value = str(row[value_col]).strip()

        if key and key.lower() != "nan":
            data[key] = value

    return data


def replace_placeholders_in_presentation(presentation: Presentation, data: dict):
    if not data:
        return

    for slide in presentation.slides:
        for shape in slide.shapes:

            if shape.has_text_frame:
                replace_placeholders_in_text_frame(shape.text_frame, data)

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_placeholders_in_text_frame(cell.text_frame, data)


def replace_placeholders_in_text_frame(text_frame, data: dict):
    for paragraph in text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)

        if not full_text:
            continue

        new_text = full_text

        for key, value in data.items():
            new_text = new_text.replace(f"{{{{{key}}}}}", str(value))

        if new_text != full_text:
            for run in paragraph.runs:
                run.text = ""

            if paragraph.runs:
                paragraph.runs[0].text = new_text
            else:
                paragraph.add_run().text = new_text


# ------------------------------------------------------------
# Diagramm-Erkennung aus PowerPoint
# ------------------------------------------------------------

def get_chart_categories(chart) -> list:
    """
    Liest die vorhandenen X-Achsen-Kategorien aus dem PowerPoint-Diagramm.
    """
    try:
        return [str(category) for category in chart.plots[0].categories]
    except Exception:
        return []


def get_chart_series_names(chart) -> list:
    """
    Liest die Namen der vorhandenen Datenreihen.
    """
    names = []

    try:
        for series in chart.series:
            names.append(str(series.name))
    except Exception:
        pass

    return names


def find_year_column(df: pd.DataFrame):
    """
    Sucht eine Jahr-Spalte.
    """
    for col in df.columns:
        if normalize_text(col) == "jahr":
            return col

    first_col = df.columns[0]

    if df[first_col].dropna().apply(is_year_like).any():
        return first_col

    return None


# ------------------------------------------------------------
# Excel-Struktur 1: breite Struktur
# ------------------------------------------------------------

def try_build_chart_data_from_wide_excel(df: pd.DataFrame, chart_categories: list):
    normalized_columns = {
        normalize_text(col): col
        for col in df.columns
    }

    matched_columns = []

    for category in chart_categories:
        normalized_category = normalize_text(category)

        if normalized_category in normalized_columns:
            matched_columns.append(normalized_columns[normalized_category])

    if len(matched_columns) == 0:
        return None

    year_column = find_year_column(df)

    if year_column:
        first_year_value = df[year_column].dropna().iloc[0]
        series_name = str(int(float(first_year_value))) if is_year_like(first_year_value) else str(first_year_value)
    else:
        series_name = "Aktuelle Werte"

    categories = []
    values = []

    for category in chart_categories:
        normalized_category = normalize_text(category)

        if normalized_category in normalized_columns:
            excel_col = normalized_columns[normalized_category]
            numeric_values = df[excel_col].apply(to_number)
            value = numeric_values.sum()
        else:
            value = 0.0

        categories.append(category)
        values.append(float(value))

    return series_name, categories, values


# ------------------------------------------------------------
# Excel-Struktur 2: lange Struktur
# ------------------------------------------------------------

def find_best_category_column(df: pd.DataFrame, chart_categories: list):
    normalized_chart_categories = {
        normalize_text(category)
        for category in chart_categories
        if normalize_text(category)
    }

    best_column = None
    best_score = 0

    for column in df.columns:
        values = {
            normalize_text(value)
            for value in df[column].dropna().astype(str).tolist()
        }

        score = len(normalized_chart_categories.intersection(values))

        if score > best_score:
            best_score = score
            best_column = column

    if best_score == 0:
        return None

    return best_column


def find_value_columns_for_long_excel(df: pd.DataFrame, category_column: str, series_names: list) -> list:
    excluded = {category_column}

    year_column = find_year_column(df)

    if year_column:
        excluded.add(year_column)

    for series_name in series_names:
        for col in df.columns:
            if normalize_text(col) == normalize_text(series_name):
                return [col]

    insgesamt_columns = [
        col for col in df.columns
        if str(col).strip().lower().startswith("insgesamt")
    ]

    if insgesamt_columns:
        return insgesamt_columns

    numeric_columns = []

    for col in df.columns:
        if col in excluded:
            continue

        converted = df[col].apply(to_number)

        if converted.sum() != 0:
            numeric_columns.append(col)

    if not numeric_columns:
        return []

    return numeric_columns


def try_build_chart_data_from_long_excel(
    df: pd.DataFrame,
    chart_categories: list,
    series_names: list
):
    category_column = find_best_category_column(df, chart_categories)

    if not category_column:
        return None

    value_columns = find_value_columns_for_long_excel(
        df=df,
        category_column=category_column,
        series_names=series_names
    )

    if not value_columns:
        return None

    year_column = find_year_column(df)

    if year_column:
        first_year_value = df[year_column].dropna().iloc[0]
        series_name = str(int(float(first_year_value))) if is_year_like(first_year_value) else str(first_year_value)
    elif series_names:
        series_name = series_names[0]
    else:
        series_name = "Aktuelle Werte"

    working_df = df.copy()
    working_df["_normalized_category"] = working_df[category_column].apply(normalize_text)

    for col in value_columns:
        working_df[col] = working_df[col].apply(to_number)

    working_df["_chart_value"] = working_df[value_columns].sum(axis=1)

    grouped = working_df.groupby("_normalized_category")["_chart_value"].sum().to_dict()

    categories = []
    values = []

    for category in chart_categories:
        normalized_category = normalize_text(category)

        categories.append(category)
        values.append(float(grouped.get(normalized_category, 0.0)))

    return series_name, categories, values


# ------------------------------------------------------------
# Hauptfunktion: Diagramme in PowerPoint ersetzen
# ------------------------------------------------------------

def replace_charts_with_matching_excel_data(
    presentation: Presentation,
    excel_path: Path
) -> int:
    df = pd.read_excel(excel_path)

    if df.empty:
        raise ValueError("Die Excel-Datei ist leer.")

    df = normalize_columns(df)

    updated_charts = 0
    found_charts = 0
    errors = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        shapes = list(slide.shapes)

        for shape_index, shape in enumerate(shapes, start=1):

            if not shape.has_chart:
                continue

            found_charts += 1

            chart = shape.chart
            chart_categories = get_chart_categories(chart)
            series_names = get_chart_series_names(chart)

            if not chart_categories:
                errors.append(f"Folie {slide_index}: Diagramm hat keine auslesbaren Kategorien.")
                continue

            chart_data_result = try_build_chart_data_from_wide_excel(
                df=df,
                chart_categories=chart_categories
            )

            if chart_data_result is None:
                chart_data_result = try_build_chart_data_from_long_excel(
                    df=df,
                    chart_categories=chart_categories,
                    series_names=series_names
                )

            if chart_data_result is None:
                errors.append(
                    f"Folie {slide_index}: Keine passenden Excel-Daten für das Diagramm gefunden."
                )
                continue

            series_name, categories, values = chart_data_result

            chart_data = CategoryChartData()
            chart_data.categories = categories
            chart_data.add_series(series_name, values)

            replace_or_recreate_chart(
                slide=slide,
                shape=shape,
                chart_data=chart_data
            )

            updated_charts += 1

    if found_charts > 0 and updated_charts == 0:
        error_text = "Es wurde ein Diagramm gefunden, aber es konnte nicht aktualisiert werden."

        if errors:
            error_text += " Details: " + " | ".join(errors)

        raise ValueError(error_text)

    return updated_charts


# ------------------------------------------------------------
# Leichte Datenvorschau ohne LibreOffice
# ------------------------------------------------------------

def extract_preview_data_from_excel(excel_path: Path):
    df = pd.read_excel(excel_path)

    if df.empty:
        return [], "Aktuelles Jahr"

    df = normalize_columns(df)
    year_column = find_year_column(df)

    if year_column:
        year_value = df[year_column].dropna().iloc[0]
        year_label = str(int(float(year_value))) if is_year_like(year_value) else str(year_value)
    else:
        year_label = "Aktuelles Jahr"

    likely_category_col = None

    for col in df.columns:
        if normalize_text(col) in ["todesursache", "ursache", "kategorie", "category"]:
            likely_category_col = col
            break

    if likely_category_col:
        value_columns = find_value_columns_for_long_excel(
            df=df,
            category_column=likely_category_col,
            series_names=[]
        )

        if value_columns:
            working_df = df.copy()

            for col in value_columns:
                working_df[col] = working_df[col].apply(to_number)

            working_df["_value"] = working_df[value_columns].sum(axis=1)

            grouped = (
                working_df
                .groupby(likely_category_col)["_value"]
                .sum()
                .reset_index()
                .rename(columns={likely_category_col: "category", "_value": "value"})
            )

            data = [
                {"category": str(row["category"]), "value": float(row["value"])}
                for _, row in grouped.iterrows()
                if float(row["value"]) > 0
            ]

            data = sorted(data, key=lambda x: x["value"], reverse=True)
            return data, year_label

    data = []

    for col in df.columns:
        if col == year_column:
            continue

        value = df[col].apply(to_number).sum()

        if value > 0:
            data.append({
                "category": str(col),
                "value": float(value)
            })

    data = sorted(data, key=lambda x: x["value"], reverse=True)
    return data, year_label


def create_excel_chart_preview_images(
    excel_path: Path,
    preview_folder: Path,
    job_id: str
) -> list:
    data, year_label = extract_preview_data_from_excel(excel_path)

    if not data:
        return []

    preview_images = []

    top10 = data[:10]

    if top10:
        image_path = preview_folder / "preview_top10.png"

        categories = [wrap_label(item["category"], 38) for item in reversed(top10)]
        values = [item["value"] for item in reversed(top10)]

        plt.figure(figsize=(18, 10))
        plt.barh(categories, values)
        plt.title(f"Top 10 Todesursachen {year_label}", fontsize=20)
        plt.xlabel("Anzahl", fontsize=14)
        plt.xticks(fontsize=11)
        plt.yticks(fontsize=10)
        plt.tight_layout()
        plt.savefig(image_path, dpi=150)
        plt.close()

        preview_images.append(
            f"{BASE_URL}/previews/{job_id}/{image_path.name}"
        )

    top5 = data[:5]
    other_value = sum(item["value"] for item in data[5:])

    if top5:
        image_path = preview_folder / "preview_top5_anteil.png"

        labels = [wrap_label(item["category"], 28) for item in top5]
        values = [item["value"] for item in top5]

        if other_value > 0:
            labels.append("Sonstige")
            values.append(other_value)

        plt.figure(figsize=(14, 10))
        plt.pie(values, labels=labels, autopct="%1.1f%%", startangle=90)
        plt.title(f"Anteil Top 5 Todesursachen {year_label}", fontsize=20)
        plt.tight_layout()
        plt.savefig(image_path, dpi=150)
        plt.close()

        preview_images.append(
            f"{BASE_URL}/previews/{job_id}/{image_path.name}"
        )

    top20 = data[:20]

    if top20:
        image_path = preview_folder / "preview_top20.png"

        categories = [wrap_label(item["category"], 42) for item in reversed(top20)]
        values = [item["value"] for item in reversed(top20)]

        plt.figure(figsize=(18, 14))
        plt.barh(categories, values)
        plt.title(f"Top 20 Todesursachen {year_label}", fontsize=20)
        plt.xlabel("Anzahl", fontsize=14)
        plt.xticks(fontsize=11)
        plt.yticks(fontsize=9)
        plt.tight_layout()
        plt.savefig(image_path, dpi=150)
        plt.close()

        preview_images.append(
            f"{BASE_URL}/previews/{job_id}/{image_path.name}"
        )

    return preview_images


def create_status_preview_images(
    presentation: Presentation,
    preview_folder: Path,
    job_id: str
) -> list:
    preview_images = []
    slide_count = len(presentation.slides)

    for index in range(slide_count):
        image = Image.new("RGB", (1920, 1080), color=(245, 246, 250))
        draw = ImageDraw.Draw(image)

        title = f"Folie {index + 1}"
        subtitle = "PowerPoint wurde erstellt"

        try:
            font_title = ImageFont.truetype("arial.ttf", 120)
            font_subtitle = ImageFont.truetype("arial.ttf", 54)
        except Exception:
            font_title = ImageFont.load_default()
            font_subtitle = ImageFont.load_default()

        draw.text((690, 430), title, fill=(44, 62, 80), font=font_title)
        draw.text((610, 570), subtitle, fill=(127, 140, 141), font=font_subtitle)

        image_path = preview_folder / f"slide_{index + 1}.png"
        image.save(image_path)

        preview_images.append(
            f"{BASE_URL}/previews/{job_id}/{image_path.name}"
        )

    return preview_images


# ------------------------------------------------------------
# Lokaler Start ohne Docker
# ------------------------------------------------------------

if __name__ == "__main__":
    import uvicorn

    print("--------------------------------------------------")
    print("PowerPoint Automation Backend startet")
    print(f"URL:        {BASE_URL}")
    print(f"App-Ordner: {APP_DIR}")
    print(f"Uploads:    {UPLOAD_DIR}")
    print(f"Outputs:    {OUTPUT_DIR}")
    print(f"Previews:   {PREVIEW_DIR}")
    print(f"Mappings:   {MAPPING_DIR}")
    print("--------------------------------------------------")

    uvicorn.run(
        app,
        host=HOST,
        port=PORT,
        reload=False
    )